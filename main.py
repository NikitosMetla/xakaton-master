# -*- coding: utf-8 -*-
import os
import json
import threading
import time
import datetime
import requests
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters import Text
from aiogram.dispatcher.filters.state import StatesGroup, State
from aiogram.types import CallbackQuery
from matplotlib import pyplot as plt
from Settings import BOT_TOKEN, MainMenu, AdminMenu, GoToQuestions, AdminStopKeyboard, MainAdmin, MainAdminMenu
from aiogram import Bot, Dispatcher, executor, types

bot = Bot(token=BOT_TOKEN, parse_mode="html")
storage = MemoryStorage()
dp = Dispatcher(bot, storage=storage)
sticker_id = 'CAACAgIAAxkBAAEHPahjwbzlOsXIHVlQguOW4s9bPK9-sAACCAADa-18CjWBoH9uCkN_LQQ'

with open("Questions.json", 'r') as file_questions:
    file_questions = json.load(file_questions)

with open('Users.json', 'r') as file_users:
    file_users = json.load(file_users)

with open("AdminList.json", 'r') as file_admins:
    file_admins = json.load(file_admins)

file_questions_index = {admin_id: 0 for admin_id in file_admins}

class InputMessage(StatesGroup):
    question = State()
    question_index_to_add = State()
    question_index_to_change = State()
    new_question = State()
    question_index_to_remove = State()
    question_to_remove_apply = State()
    questions_start = State()
    user_telegram_id = State()

def create_histogram() -> None:
    requests_by_the_hour = [0 for _ in range(24)]
    timestamp = int(time.time())
    for TelegramId in file_users:
        if file_users[TelegramId][1]["TimeEnd"] is not None and datetime.datetime.fromtimestamp(
                file_users[TelegramId][1]["TimeEnd"]).strftime('%d') == datetime.datetime.fromtimestamp(
                timestamp).strftime('%d'):
            requests_by_the_hour[
                int(datetime.datetime.fromtimestamp(file_users[TelegramId][1]["TimeEnd"]).strftime('%H'))] += 1
    x_list = []
    for hour in range(24):
        hour = str(hour)
        if len(hour) == 1:
            hour = '0' + hour
        x_list.append(hour)
    y_list = requests_by_the_hour

    plt.title('Количество заявок по часам')
    plt.xlabel('Часы')
    plt.ylabel('Количество заявок')

    plt.bar(x_list, y_list)
    plt.savefig('Histogram.png', dpi='figure',
                bbox_inches=None, pad_inches=1,
                facecolor='auto', edgecolor='auto',
                backend=None)
    return None

def create_pptx(TelegramId: str) -> None:
    root = Presentation('Шаблон.pptx')
    for index_question in range(7):
        slide = root.slides[index_question]
        if index_question == 0:
            if file_users[TelegramId][0][index_question]["Answer"]["Icon"] is not None:
                file_path = requests.get(
                    f'https://api.telegram.org/bot{BOT_TOKEN}/getFile?file_id={file_users[TelegramId][0][index_question]["Answer"]["Icon"]["file_id"]}').text
                file_path = json.loads(file_path)
                if file_path['result']:
                    file_path = str(file_path['result']['file_path'])
                    icon = requests.get(f'https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}').content
                    img_path = TelegramId + '_' + file_path.split("/")[1]
                    with open(img_path, 'wb') as new_file:
                        new_file.write(icon)
                    im = Image.open(img_path)
                    height, width = normalize_size(im.height, im.width)
                    left = Inches(5.5)
                    top = Inches(4.1)
                    slide.shapes.add_picture(img_path, left, top, width=Inches(width / 100),
                                             height=Inches(height / 100))
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
        elif index_question == 1:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
        elif index_question == 2:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
            if file_users[TelegramId][0][index_question]["Answer"]["Icon"] is not None:
                file_path = requests.get(
                    f'https://api.telegram.org/bot{BOT_TOKEN}/getFile?file_id={file_users[TelegramId][0][index_question]["Answer"]["Icon"]["file_id"]}').text
                file_path = json.loads(file_path)
                if file_path['result']:
                    file_path = str(file_path['result']['file_path'])
                    icon = requests.get(f'https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}').content
                    img_path = TelegramId + '_' + file_path.split("/")[1]
                    with open(img_path, 'wb') as new_file:
                        new_file.write(icon)
                    im = Image.open(img_path)
                    height, width = normalize_size(im.height, im.width)
                    left = Inches(5.5)
                    top = Inches(0.15)
                    slide.shapes.add_picture(img_path, left, top, width=Inches(width / 100),
                                             height=Inches(height / 100))
        elif index_question == 3:
            incomes = str(file_users[TelegramId][0][index_question]["Answer"]["Text"]).split(' ')
            real_incomes = [int(coin) for coin in incomes if coin.isdigit()][:3]
            chart_data = CategoryChartData()
            chart_data.categories = ['2020', '2021', '2022']
            chart_data.add_series('Годовая прибыль', tuple(real_incomes))
            x, y, cx, cy = Inches(4), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        elif index_question == 4:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
        if file_users[TelegramId][0][index_question]["Answer"]["Icon"] is not None:
            file_path = requests.get(
                f'https://api.telegram.org/bot{BOT_TOKEN}/getFile?file_id={file_users[TelegramId][0][index_question]["Answer"]["Icon"]["file_id"]}').text
            file_path = json.loads(file_path)
            if file_path['result']:
                file_path = str(file_path['result']['file_path'])
                icon = requests.get(f'https://api.telegram.org/file/bot{BOT_TOKEN}/{file_path}').content
                img_path = TelegramId + '_' + file_path.split("/")[1]
                with open(img_path, 'wb') as new_file:
                    new_file.write(icon)
                im = Image.open(img_path)
                height, width = normalize_size(im.height, im.width)
                left = Inches(5.5)
                top = Inches(4.1)
                slide.shapes.add_picture(img_path, left, top, width=Inches(width / 100),
                                         height=Inches(height / 100))
        elif index_question == 5:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
        elif index_question == 6:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = str(file_users[TelegramId][0][index_question]["Answer"]["Text"])
    root.save(f'{TelegramId}.pptx')
    return None

def create_docx(TelegramId: str) -> None:
    document = Document()
    document.add_heading('Контент для Мемо', 0)
    for id_question, question in enumerate(file_users[TelegramId][0]):
        p = document.add_paragraph(
            f'Вопрос №{id_question + 1}: ', style='List Bullet'
        )
        p.add_run(f" {file_users[TelegramId][0][id_question]['Question']}").italic = True
        p.add_run(
            f"\nОтвет на вопрос: {file_users[TelegramId][0][id_question]['Answer']['Text']}").italic = True
    document.save(f'{TelegramId}.docx')
    return None

def clear_temp_files(TelegramId: str) -> None:
    TelegramId = str(TelegramId)
    for file in os.listdir():
        if TelegramId in file or 'Histogram.png' in file:
            try:
                os.remove(file)
            except ExceptionGroup:
                return clear_temp_files(TelegramId)

def normalize_size(height: int, width: int) -> tuple:
    max_height = 285
    max_width = 427
    while width > max_width or height > max_height:
        if width > max_width:
            const = max_width / width
            width = max_width
            height = height * const
        elif height > max_height:
            const = max_height / height
            height = max_height
            width = width * const
    return height, width


@dp.message_handler(commands=['start'])
async def start(message: types.Message):
    await message.answer_sticker(sticker_id)
    await message.answer('<b>Приветствуем тебя в нашем боте!\nЗдесь ты cможешь отправить <u>анкету</u> со своими предложением.\nАнкетирование будет проходить <u>в формате диалога</u> с ботом!</b>', reply_markup=MainMenu)


@dp.message_handler(lambda message: message.text == "Заполнить заявку")
async def fill_application(message: types.Message):
    await message.answer("<b>Сейчас вам будет задан ряд вопросов. Вы готовы ответить на них?</b>", reply_markup=GoToQuestions)


@dp.message_handler(lambda message: message.text == "Справка")
async def fill_application(message: types.Message):
    await message.answer("<b>Данный телеграм-бот предназначен для отправки своих предложений по проектам для дальнейшего инвестирования.  Бот имеет функционал как для работы с обычными пользователями, так и с аналитиками, обрабатывающими анкеты.</b>")


@dp.callback_query_handler(Text("Yes"))
async def start_questions(call: CallbackQuery):
    await bot.delete_message(call.from_user.id, call.message.message_id)
    if len(file_questions['Questions']) != 0:
        file_users[str(call.from_user.id)] = [[], {"TimeStart": None, "TimeEnd": None}]
        answers_list = []
        for id_question, question in enumerate(file_questions['Questions']):
            answers_list.append({"Question": question, "Answer": {"Icon": None, "Text": None}})
        file_users[str(call.from_user.id)][0] = answers_list
        file_users[str(call.from_user.id)][1]["TimeStart"] = int(time.time())
        with open("Users.json", 'w') as file_users_two:
            json.dump(file_users, file_users_two, indent=2)
        await call.message.answer('<b>Постарайтесь как можно точнее ответить на вопросы, нам будет интересно услышать Вас!</b>\n<i>Если не знаете, что ответить на вопрос,то отправляйте любой символ</i>')
        await call.message.answer(f"<b>Вопрос №1</b>: <i>{file_questions['Questions'][0]}</i>")
        await InputMessage.questions_start.set()
    else:
        await call.message.answer('<b>Подождите, когда администрация добавит вопросы</b>')


@dp.message_handler(content_types=['photo', 'text'], state=InputMessage.questions_start)
async def send_question_to_user(message: types.Message, state: FSMContext):
    id_question = 0
    for id_question in range(len(file_users[str(message.from_user.id)][0])):
        if file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] is None and file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] is None:
            file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] = message.caption if message.caption is not None else message.text
            file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] = {"file_id": message.photo[-1]["file_id"], "file_unique_id": message.photo[-1]["file_unique_id"]} if len(message.photo) != 0 else None
            break
    if id_question == len(file_users[str(message.from_user.id)][0]) - 1:
        file_users[str(message.from_user.id)][1]['TimeEnd'] = int(time.time())
        with open("Users.json", 'w') as file_users_two:
            json.dump(file_users, file_users_two, indent=2)
        await message.answer('<b>Вы успешно прошли опрос. Ваша анкета сохранена и отправлена!</b>')
        await state.finish()
        """СОЗДАНИЕ DOCX"""
        create_docx(str(message.from_user.id))
        create_pptx(str(message.from_user.id))
        with open(f'{message.from_user.id}.docx', 'rb') as document_path:
            with open(f'{message.from_user.id}.pptx', 'rb') as pptx_path:
                timestamp = int(time.time())
                min_accepted = []
                for IdAdmin in file_admins:
                    if file_admins[IdAdmin]['LastAccepted'] is None:
                        break
                    else:
                        min_accepted.append([file_admins[IdAdmin]["LastAccepted"] - timestamp, IdAdmin])
                else:
                    IdAdmin = sorted(min_accepted)[0][1]
                await bot.send_message(IdAdmin, '<b>Вам пришла новая анкета!</b>')
                await bot.send_document(chat_id=IdAdmin, document=document_path)
                await bot.send_document(chat_id=IdAdmin, document=pptx_path)
                file_admins[IdAdmin]['LastAccepted'] = int(time.time())
                file_admins[IdAdmin]['Accepted'] += 1
                with open('AdminList.json', 'w') as file_admins_two:
                    json.dump(file_admins, file_admins_two, indent=2)
            """ОЧИСТКА ВРЕМЕННЫХ ФАЙЛОВ"""
            threading.Thread(target=clear_temp_files, args=(message.from_user.id,)).start()
            with open('AdminList.json', 'w') as file_admins_two:
                json.dump(file_admins, file_admins_two, indent=2)
    else:
        for id_question in range(len(file_users[str(message.from_user.id)][0])):
            if file_users[str(message.from_user.id)][0][id_question]["Answer"]["Text"] is None and file_users[str(message.from_user.id)][0][id_question]["Answer"]["Icon"] is None:
                await message.answer(f"<b>Вопрос №{id_question + 1}:</b> <i>{file_users[str(message.from_user.id)][0][id_question]['Question']}</i>")
                break


@dp.callback_query_handler(Text(["No"]))
async def cansel_questions(call: CallbackQuery):
    await bot.delete_message(call.from_user.id, call.message.message_id)
    await call.message.answer('Выберите свои дальнейшие действия', reply_markup=MainMenu)


@dp.message_handler(commands=['admin'])
async def admin_settings(message: types.Message):
    if str(message.from_user.id) in file_admins:
        if str(message.from_user.id) in MainAdmin:
            await message.answer('<b>Успешный вход в админ-панель!</b>', reply_markup=MainAdminMenu)
            create_histogram()
            with open('Histogram.png', 'rb') as file:
                await message.answer('<b>Статистика заявок за день:</b>')
                await bot.send_photo(message.from_user.id, photo=file)
        else:
            await message.answer('<b>Успешный вход в админ-панель!</b>', reply_markup=AdminMenu)
    else:
        await message.answer('<b>Вы не являетесь администратором!</b>')


@dp.message_handler(lambda message: message.text == "Добавить/Удалить админа")
async def add_remove_admin(message: types.Message):
    if str(message.from_user.id) in MainAdmin:
        await message.answer('<b>Напишите TelegramId человека, кому вы хотите предоставить права администратора:</b>', reply_markup=AdminStopKeyboard)
        await InputMessage.user_telegram_id.set()
    else:
        await message.answer('<b>У вас нет прав на выполнение этой команды</b>')


@dp.message_handler(lambda message: message.text == 'Список админов')
async def list_of_admins(message: types.Message):
    if str(message.from_user.id) in MainAdmin:
        admins_list = '<b>Список администраторов:</b>\n\n'
        for admin in file_admins:
            if file_admins[admin]["LastAccepted"] is None:
                admins_list += f'<b>TelegramId: <code>{admin}</code>\nПросмотрено заявок: <i>{file_admins[admin]["Accepted"]}</i>\nДата последней принятой заявки: <i>{None}</i></b>\n\n'
            else:
                admins_list += f'<b>TelegramId: <code>{admin}</code>\nПросмотрено заявок: <i>{file_admins[admin]["Accepted"]}</i>\nДата последней принятой заявки: <i>{datetime.datetime.utcfromtimestamp(int(file_admins[admin]["LastAccepted"])).strftime("%Y-%m-%d %H:%M:%S")}</i></b>\n\n'
        await message.answer(admins_list, reply_markup=MainAdminMenu)
    else:
        await message.answer('<b>У вас нет прав на выполнение этой команды</b>')


@dp.message_handler(state=InputMessage.user_telegram_id)
async def check_to_add_or_remove_user_telegram_id(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    else:
        if str(message.text).isdigit():
            if message.text in file_admins and message.text not in MainAdmin:
                del (file_admins[str(message.text)])
                await message.answer('<b>Администратор был удалён</b>', reply_markup=MainAdminMenu)
                with open('AdminList.json', 'w') as file_admins_two:
                    json.dump(file_admins, file_admins_two, indent=2, ensure_ascii=False)
            elif message.text in MainAdmin:
                await message.answer('<b>Вы не можете удалить главного администратора через данную панель</b>', MainAdminMenu)
            else:
                file_admins[str(message.text)] = {"Accepted": 0, "LastAccepted": None}
                await message.answer('<b>Администратор был добавлен</b>', reply_markup=MainAdminMenu)
                with open('AdminList.json', 'w') as file_admins_two:
                    json.dump(file_admins, file_admins_two, indent=2, ensure_ascii=False)
        else:
            await message.answer('<b>Вы неправильно ввели TelegramId</b>')
    await state.finish()


@dp.message_handler(lambda message: message.text == "Главное меню")
async def out_admin_settings(message: types.Message):
    await message.answer('<b>Выберите свои дальнейшие действия:</b>', reply_markup=MainMenu)


@dp.message_handler(lambda message: message.text == "Список вопросов")
async def send_questions_to_admin(message: types.Message):
    if len(file_questions["Questions"]) != 0:
        questions = '<b>Первые 7 вопросов <u>нежелательно менять</u>, на основе них строится презентации</b>\n\n'
        for id_question, question in enumerate(file_questions["Questions"]):
            questions += f'<b>№{id_question + 1}:</b> <i>{question}</i>\n\n'
        await message.answer(questions)
    else:
        await message.answer('<b>Список вопросов пуст</b>')


@dp.message_handler(lambda message: message.text == "Добавить вопрос")
async def send_new_question_out_admin(message: types.Message):
    await message.answer('<b>Каким по номеру должен стоять новый вопрос из всех вопросов:</b>',
                         reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_add.set()


@dp.message_handler(state=InputMessage.question_index_to_add)
async def add_new_question(message: types.Message, state: FSMContext):
    try:
        if message.text == 'Отменить действие':
            await state.finish()
            if str(message.from_user.id) not in MainAdmin:
                await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
            else:
                await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
        elif int(message.text) > 0:
            file_questions_index[str(message.from_user.id)] = int(message.text)
            await state.finish()
            await message.answer('<b>Введите ваш вопрос:</b>', reply_markup=AdminStopKeyboard)
            await InputMessage.question.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    except ExceptionGroup:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.question)
async def replace_new_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        if str(message.from_user.id) not in MainAdmin:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        else:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    else:
        id_question = file_questions_index[str(message.from_user.id)]
        questions = list(file_questions["Questions"])[:id_question - 1] + [message.text] + list(
            file_questions["Questions"])[id_question - 1:]
        file_questions["Questions"] = questions
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2, ensure_ascii=False)
        await message.answer('<b>Вопрос успешно добавлен в список!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()


@dp.message_handler(lambda message: message.text == "Изменить вопрос")
async def change_question_index_out_admin(message: types.Message):
    await message.answer('<b>Пришлите номер вопроса, какой вы хотите изменить</b>', reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_change.set()


@dp.message_handler(state=InputMessage.question_index_to_change)
async def change_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        if str(message.from_user.id) not in MainAdmin:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        else:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    if str(message.text).isdigit():
        if 0 < int(message.text) <= len(file_questions['Questions']):
            file_questions_index[str(message.from_user.id)] = int(message.text)
            await state.finish()
            await message.answer('Введите ваш новый вопрос', reply_markup=AdminStopKeyboard)
            await InputMessage.new_question.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    else:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.new_question)
async def replace_new_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        if str(message.from_user.id) not in MainAdmin:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        else:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    else:
        id_question = file_questions_index[str(message.from_user.id)]
        file_questions["Questions"][id_question - 1] = message.text
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2, ensure_ascii=False)
        await message.answer('<b>Вопрос успешно изменён!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()


@dp.message_handler(lambda message: message.text == "Удалить вопрос")
async def send_question_index_to_remove_out_admin(message: types.Message):
    await message.answer('<b>Пришлите номер вопроса, какой вы хотите удалить</b>', reply_markup=AdminStopKeyboard)
    await InputMessage.question_index_to_remove.set()


@dp.message_handler(state=InputMessage.question_index_to_remove)
async def remove_question_apply(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие':
        await state.finish()
        if str(message.from_user.id) not in MainAdmin:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        else:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    if str(message.text).isdigit():
        if 0 < int(message.text) <= len(file_questions['Questions']):
            file_questions_index[str(message.from_user.id)] = int(message.text)
            await state.finish()
            await message.answer(
                f'<b>Вы точно хотите удалить вопрос:</b> <i>{file_questions["Questions"][int(message.text) - 1][:35]}</i>..?',
                reply_markup=AdminStopKeyboard)
            await message.answer('<b>Ответьте: <i>Да / Нет</i></b>', reply_markup=AdminStopKeyboard)
            await InputMessage.question_to_remove_apply.set()
        else:
            await message.answer('<b>Введите правильно индекс</b>')
    else:
        await message.answer('<b>Введите правильно индекс</b>')


@dp.message_handler(state=InputMessage.question_to_remove_apply)
async def remove_question(message: types.Message, state: FSMContext):
    if message.text == 'Отменить действие' or message.text == 'Нет':
        await state.finish()
        if str(message.from_user.id) not in MainAdmin:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=AdminMenu)
        else:
            await message.answer('<b>Действие успешно отменено</b>', reply_markup=MainAdminMenu)
    elif message.text == 'Да':
        id_question = file_questions_index[str(message.from_user.id)]
        del (file_questions["Questions"][id_question - 1])
        with open('Questions.json', 'w') as file_questions_two:
            json.dump(file_questions, file_questions_two, indent=2, ensure_ascii=False)
        await message.answer('<b>Вопрос успешно удалён!\nСейчас список вопросов выглядит так:</b>',
                             reply_markup=AdminMenu)
        await send_questions_to_admin(message)
        await state.finish()
    else:
        await message.answer('<b>Ответьте: <i>Да / Нет</i></b>', reply_markup=AdminStopKeyboard)


@dp.message_handler(content_types=["photo", "text"])
async def incomprehensible_message(message: types.Message):
    await message.answer('<b>Я вас не понимаю :(</b>')

if __name__ == "__main__":
    executor.start_polling(dp, skip_updates=True, timeout=None)
